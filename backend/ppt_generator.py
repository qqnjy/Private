"""Build the weekly social-media report PPT using the IGS master template.

Deck (FB + IG only; Threads explicitly excluded):
  1. Cover                                     (layout 0, white)
  2. 一、整體表現總結                            (layout 2, white)
  3. Facebook 成效對比 + 一週回顧                (layout 3, custom table + bullets)
  4. Facebook TOP 5 貼文排行（含縮圖）            (layout 3, custom table)
  5. Instagram 成效對比 + 一週回顧                (layout 3, custom table + bullets)
  6. Instagram TOP 5 貼文排行                    (layout 3, custom table)
  7. 三、後續規劃                                 (layout 2)
  8. Thank You                                  (layout 3, title only)
"""
import io
import os
import httpx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "templates", "igs_master.pptx")

LAYOUT_COVER = 0
LAYOUT_CONTENT = 2
LAYOUT_BLANK = 3   # title-only — we add tables manually onto these slides

PH_TITLE = 0
PH_SUBTITLE = 1
PH_BODY1 = 11
PH_BODY2 = 12
PH_PICTURE = 10

BRAND_RED = RGBColor(0xE5, 0x01, 0x12)
BRAND_GREEN = RGBColor(0x00, 0xA3, 0x4E)
HEADER_GOLD = RGBColor(0xF5, 0xA8, 0x23)
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)
TEXT_MUTED = RGBColor(0x7F, 0x8C, 0x8D)
ROW_ALT = RGBColor(0xF8, 0xF9, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GROWTH_UP = RGBColor(0xE7, 0x4C, 0x3C)   # red — Chinese convention: up = red
GROWTH_DOWN = RGBColor(0x27, 0xAE, 0x60)  # green — down = green


# ---------- helpers ----------

def _clear_existing_slides(prs):
    sld_id_lst = prs.slides._sldIdLst
    for sld in list(sld_id_lst):
        rId = sld.rId
        sld_id_lst.remove(sld)
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


def _ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _remove_placeholder(slide, idx):
    ph = _ph(slide, idx)
    if ph is not None:
        sp = ph._element
        sp.getparent().remove(sp)


def _set_text(slide, idx, text):
    ph = _ph(slide, idx)
    if ph is not None:
        ph.text = text or ""
        return ph
    return None


def _fill_body(placeholder, paragraphs):
    if placeholder is None:
        return
    tf = placeholder.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(paragraphs):
        text = item.get("text", "")
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if "level" in item:
            p.level = item["level"]
        r = p.add_run()
        r.text = text
        if "size" in item:
            r.font.size = Pt(item["size"])
        if item.get("bold"):
            r.font.bold = True
        if "color" in item:
            r.font.color.rgb = item["color"]


def _download_image(url):
    if not url:
        return None
    try:
        r = httpx.get(url, timeout=10.0, follow_redirects=True)
        r.raise_for_status()
        return io.BytesIO(r.content)
    except Exception:
        return None


def _format_num(n):
    if n is None:
        return "—"
    if isinstance(n, float):
        if n.is_integer():
            n = int(n)
        else:
            return f"{n:,.1f}"
    return f"{n:,}"


def _pct_change(this_val, last_val):
    if not last_val:
        return None
    if this_val is None:
        return None
    return ((this_val - last_val) / last_val) * 100


def _format_pct(pct):
    if pct is None:
        return ""
    sign = "+" if pct >= 0 else ""
    return f"{sign}{pct:.1f}%"


# ---------- text cell styling ----------

def _derive_post_type(p: dict) -> str:
    """Fallback: derive post type from media_type / status_type if post_type missing."""
    if p.get("post_type"):
        return p["post_type"]
    mt = (p.get("media_type") or "").lower()
    if mt in ("video", "reels"):
        return "梗影"
    if mt == "album" or mt == "carousel_album":
        return "圖文"
    if mt in ("photo", "image"):
        return "梗圖"
    return "—"


def _add_slide_title(slide, text, color=BRAND_RED):
    """Add a top-left title text box (since LAYOUT_BLANK's title is centered)."""
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.30), Inches(9.0), Inches(0.65))
    tf = tb.text_frame
    tf.margin_left = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = color


def _set_cell(cell, text, *, size=11, bold=False, color=TEXT_DARK, bg=None, align=PP_ALIGN.CENTER):
    cell.text = ""
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


# ---------- slide builders ----------

def _add_cover(prs, brand_name, week_range):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_COVER])
    _set_text(slide, PH_TITLE, f"{brand_name} 社群週報")
    _set_text(slide, PH_SUBTITLE, week_range)


def _add_summary_slide(prs, slides_data):
    summary = (slides_data or {}).get("summary", {}) or {}
    body1 = []
    body2 = []
    if summary.get("fb"):
        body1.append({"text": "Facebook", "bold": True, "size": 18, "color": BRAND_RED})
        body1.append({"text": summary["fb"], "size": 14, "color": TEXT_DARK})
    if summary.get("ig"):
        body2.append({"text": "Instagram", "bold": True, "size": 18, "color": BRAND_RED})
        body2.append({"text": summary["ig"], "size": 14, "color": TEXT_DARK})

    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    _set_text(slide, PH_TITLE, "一、整體表現總結")
    _remove_placeholder(slide, PH_SUBTITLE)
    if body1:
        _fill_body(_ph(slide, PH_BODY1), body1)
    else:
        _remove_placeholder(slide, PH_BODY1)
    if body2:
        _fill_body(_ph(slide, PH_BODY2), body2)
    else:
        _remove_placeholder(slide, PH_BODY2)
    _remove_placeholder(slide, PH_PICTURE)


def _add_comparison_slide(prs, platform_label, week_range, this_week, last_week, ai_info, post_count_this, post_count_last):
    """Slide: comparison table on left, AI bullets on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _remove_placeholder(slide, PH_TITLE)  # built-in title is centered low; we add our own
    _add_slide_title(slide, f"{platform_label}　成效對比")

    # ----- LEFT: comparison table -----
    rows_data = [
        ("貼文總數", post_count_this, post_count_last),
        ("總觀看數", this_week.get("total_views"), last_week.get("total_views") if last_week else None),
        ("總互動數", this_week.get("total_interactions"), last_week.get("total_interactions") if last_week else None),
        ("總按讚數", this_week.get("total_likes"), last_week.get("total_likes") if last_week else None),
        ("總留言數", this_week.get("total_comments"), last_week.get("total_comments") if last_week else None),
        ("總分享數", this_week.get("total_shares"), last_week.get("total_shares") if last_week else None),
        ("平均觀看", this_week.get("avg_view"), last_week.get("avg_view") if last_week else None),
        ("互動率(%)", this_week.get("interaction_rate"), last_week.get("interaction_rate") if last_week else None),
    ]

    rows = len(rows_data) + 1
    cols = 4
    left = Inches(0.45)
    top = Inches(1.30)
    width = Inches(4.55)
    height = Inches(3.80)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table
    # Column widths
    tbl.columns[0].width = Inches(1.20)
    tbl.columns[1].width = Inches(1.10)
    tbl.columns[2].width = Inches(1.10)
    tbl.columns[3].width = Inches(1.15)

    # Header row
    headers = ["指標", "本週", "上週", "變化"]
    for c, h in enumerate(headers):
        _set_cell(tbl.cell(0, c), h, size=12, bold=True, color=WHITE, bg=HEADER_GOLD)

    for ri, (label, this_v, last_v) in enumerate(rows_data, start=1):
        bg = ROW_ALT if ri % 2 == 0 else WHITE
        _set_cell(tbl.cell(ri, 0), label, size=11, bold=True, color=TEXT_DARK, bg=bg, align=PP_ALIGN.LEFT)
        _set_cell(tbl.cell(ri, 1), _format_num(this_v), size=11, color=TEXT_DARK, bg=bg)
        _set_cell(tbl.cell(ri, 2), _format_num(last_v) if last_v is not None else "—", size=11, color=TEXT_MUTED, bg=bg)
        pct = _pct_change(this_v, last_v) if last_v is not None else None
        if pct is None:
            _set_cell(tbl.cell(ri, 3), "—", size=11, color=TEXT_MUTED, bg=bg)
        else:
            color = GROWTH_UP if pct > 0 else (GROWTH_DOWN if pct < 0 else TEXT_DARK)
            _set_cell(tbl.cell(ri, 3), _format_pct(pct), size=11, bold=True, color=color, bg=bg)

    # ----- RIGHT: AI bullets -----
    ai_info = ai_info or {}
    overview = ai_info.get("overview") or []
    top_posts = ai_info.get("top_posts") or []

    tb_right = slide.shapes.add_textbox(Inches(5.20), Inches(1.30), Inches(4.50), Inches(3.95))
    tf = tb_right.text_frame
    tf.word_wrap = True

    def _line(text, size=12, bold=False, color=TEXT_DARK, first=False):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

    _line("本週概況", size=14, bold=True, color=BRAND_RED, first=True)
    if overview:
        for o in overview:
            _line(f"• {o}", size=11)
    else:
        _line("• 本週數據自動結算，待 AI 補充。", size=11, color=TEXT_MUTED)

    _line("", size=4)
    _line("一週回顧 — TOP 貼文", size=14, bold=True, color=BRAND_RED)
    if top_posts:
        for t in top_posts[:3]:
            head = f"TOP {t.get('rank', '?')}：{(t.get('summary') or '')[:30]}"
            _line(head, size=11, bold=True)
            if t.get("comment"):
                _line(f"   {t['comment']}", size=10, color=TEXT_DARK)
    else:
        _line("• 待 AI 補充", size=11, color=TEXT_MUTED)


def _add_top5_slide(prs, platform_label, posts, platform):
    """Slide: TOP 5 ranking table with thumbnails."""
    if not posts:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _remove_placeholder(slide, PH_TITLE)
    _add_slide_title(slide, f"{platform_label}　TOP 5 貼文")

    # Columns: 名次 | 縮圖 | 類型 | 觀看 | 觸及 | 讚 | 留言 | 分享 | 互動率
    cols_def = [
        ("名次", Inches(0.55)),
        ("縮圖", Inches(1.20)),
        ("類型", Inches(0.70)),
        ("觀看", Inches(1.00)),
        ("觸及", Inches(0.90)),
        ("讚", Inches(0.70)),
        ("留言", Inches(0.70)),
        ("分享", Inches(0.70)),
        ("互動率", Inches(0.85)),
    ]
    n_rows = min(5, len(posts)) + 2  # header + posts + 平均
    n_cols = len(cols_def)
    total_w = sum(c[1] for c in cols_def)
    left = (Inches(10.0) - total_w) // 2
    top = Inches(1.30)
    # Note: actual height grows from row content, but specify base
    height = Inches(3.80)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, height)
    tbl = table_shape.table
    for i, (_, w) in enumerate(cols_def):
        tbl.columns[i].width = w

    # Row 0 — header
    for c, (label, _) in enumerate(cols_def):
        _set_cell(tbl.cell(0, c), label, size=11, bold=True, color=WHITE, bg=HEADER_GOLD)

    # Set the row holding images to be taller
    tbl.rows[0].height = Inches(0.32)
    for ri in range(1, n_rows):
        tbl.rows[ri].height = Inches(0.62)

    # Rows 1..5 — posts
    use_posts = posts[:5]
    sum_views = sum_reach = sum_likes = sum_comments = sum_shares = 0
    rate_count = 0
    sum_rate = 0.0
    for idx, p in enumerate(use_posts, start=1):
        bg = ROW_ALT if idx % 2 == 0 else WHITE
        if platform == "fb":
            views = p.get("total_views") or p.get("video_views") or 0
            reach = 0  # FB reach not available
            likes = p.get("reactions") or 0
            comments = p.get("comments") or 0
            shares = p.get("shares") or 0
        else:
            views = p.get("views") or 0
            reach = p.get("reach") or 0
            likes = p.get("likes") or 0
            comments = p.get("comments") or 0
            shares = p.get("shares") or 0
        interactions = likes + comments + shares
        rate = (interactions / views * 100) if views > 0 else 0
        sum_views += views; sum_reach += reach
        sum_likes += likes; sum_comments += comments; sum_shares += shares
        if views > 0:
            sum_rate += rate
            rate_count += 1

        _set_cell(tbl.cell(idx, 0), str(idx), size=11, bold=True, color=TEXT_DARK, bg=bg)
        # Image cell: empty text, then add picture overlay
        _set_cell(tbl.cell(idx, 1), "", size=10, bg=bg)
        _set_cell(tbl.cell(idx, 2), _derive_post_type(p), size=10, color=TEXT_DARK, bg=bg)
        _set_cell(tbl.cell(idx, 3), _format_num(views), size=10, color=TEXT_DARK, bg=bg)
        _set_cell(tbl.cell(idx, 4), _format_num(reach) if reach else "—", size=10, color=TEXT_DARK, bg=bg)
        _set_cell(tbl.cell(idx, 5), _format_num(likes), size=10, color=TEXT_DARK, bg=bg)
        _set_cell(tbl.cell(idx, 6), _format_num(comments), size=10, color=TEXT_DARK, bg=bg)
        _set_cell(tbl.cell(idx, 7), _format_num(shares), size=10, color=TEXT_DARK, bg=bg)
        _set_cell(tbl.cell(idx, 8), f"{rate:.2f}%" if views > 0 else "—", size=10, color=TEXT_DARK, bg=bg)

        # Embed thumbnail INTO the image cell area
        img_stream = _download_image(p.get("media_url"))
        if img_stream:
            # Compute image cell's bounding box
            cell_x = left + sum(c[1] for c in cols_def[:1])
            row_top_offset = sum((tbl.rows[r].height for r in range(idx)), Inches(0))
            cell_y = top + row_top_offset
            img_w = cols_def[1][1] - Inches(0.10)
            img_h = tbl.rows[idx].height - Inches(0.08)
            try:
                slide.shapes.add_picture(
                    img_stream, cell_x + Inches(0.05), cell_y + Inches(0.04),
                    width=img_w, height=img_h
                )
            except Exception:
                pass

    # Last row — 平均值
    n_used = len(use_posts) or 1
    avg_row = n_rows - 1
    bg = ROW_ALT
    _set_cell(tbl.cell(avg_row, 0), "平均", size=11, bold=True, color=TEXT_DARK, bg=bg)
    _set_cell(tbl.cell(avg_row, 1), "—", size=10, bg=bg)
    _set_cell(tbl.cell(avg_row, 2), "—", size=10, bg=bg)
    _set_cell(tbl.cell(avg_row, 3), _format_num(round(sum_views / n_used)), size=10, bold=True, bg=bg)
    _set_cell(tbl.cell(avg_row, 4), _format_num(round(sum_reach / n_used)) if sum_reach else "—", size=10, bold=True, bg=bg)
    _set_cell(tbl.cell(avg_row, 5), _format_num(round(sum_likes / n_used)), size=10, bold=True, bg=bg)
    _set_cell(tbl.cell(avg_row, 6), _format_num(round(sum_comments / n_used)), size=10, bold=True, bg=bg)
    _set_cell(tbl.cell(avg_row, 7), _format_num(round(sum_shares / n_used)), size=10, bold=True, bg=bg)
    _set_cell(tbl.cell(avg_row, 8), f"{(sum_rate / rate_count):.2f}%" if rate_count else "—", size=10, bold=True, bg=bg)


def _add_plans_slide(prs, slides_data):
    plans = (slides_data or {}).get("plans") or []
    if not plans:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    _set_text(slide, PH_TITLE, "三、後續規劃")
    _remove_placeholder(slide, PH_SUBTITLE)
    _remove_placeholder(slide, PH_PICTURE)

    half = (len(plans) + 1) // 2
    body1 = []
    body2 = []
    for i, p in enumerate(plans, 1):
        head = f"{i}. {p.get('platform', '')}  {p.get('title', '')}".strip()
        target = body1 if i <= half else body2
        target.append({"text": head, "bold": True, "size": 16, "color": BRAND_RED})
        if p.get("detail"):
            target.append({"text": p["detail"], "size": 12, "color": TEXT_DARK})
        target.append({"text": "", "size": 4})

    _fill_body(_ph(slide, PH_BODY1), body1)
    if body2:
        _fill_body(_ph(slide, PH_BODY2), body2)
    else:
        _remove_placeholder(slide, PH_BODY2)


def _add_closing_slide(prs):
    # Thank-you slide: keep the layout's centered title for visual symmetry
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _set_text(slide, PH_TITLE, "Thank You")


def _load_posts_and_summary(brand_name, week_range):
    """Return (posts_dict, this_summary, last_summary)."""
    if " ~ " not in week_range:
        return {"fb": [], "ig": []}, {}, {}
    try:
        from posts_repo import get_cached_posts, get_week_summary
        from datetime import datetime, timedelta

        start_s, end_s = [s.strip() for s in week_range.split(" ~ ", 1)]
        start_d = datetime.strptime(start_s, "%Y-%m-%d").date()
        end_d = datetime.strptime(end_s, "%Y-%m-%d").date()
        prev_start = (start_d - timedelta(days=7)).isoformat()
        prev_end = (end_d - timedelta(days=7)).isoformat()

        cached = get_cached_posts(brand_name, start_s, end_s) or {"fb": [], "ig": []}
        this_summary = get_week_summary(brand_name, start_s, end_s)
        last_summary = get_week_summary(brand_name, prev_start, prev_end)
        return cached, this_summary, last_summary
    except Exception as e:
        print(f"load_posts_and_summary failed: {e}")
        return {"fb": [], "ig": []}, {}, {}


def build_ppt(brand_name: str, week_range: str, slides_data: dict) -> bytes:
    prs = Presentation(TEMPLATE_PATH)
    _clear_existing_slides(prs)

    cached, this_summary, last_summary = _load_posts_and_summary(brand_name, week_range)
    fb_posts = cached.get("fb") or []
    ig_posts = cached.get("ig") or []

    _add_cover(prs, brand_name, week_range)
    _add_summary_slide(prs, slides_data or {})

    _add_comparison_slide(
        prs, "Facebook", week_range,
        this_summary.get("fb") or {}, last_summary.get("fb") if last_summary else {},
        (slides_data or {}).get("fb") or {},
        len(fb_posts), (last_summary.get("fb") or {}).get("post_count") if last_summary else None,
    )
    _add_top5_slide(prs, "Facebook", fb_posts, "fb")

    _add_comparison_slide(
        prs, "Instagram", week_range,
        this_summary.get("ig") or {}, last_summary.get("ig") if last_summary else {},
        (slides_data or {}).get("ig") or {},
        len(ig_posts), (last_summary.get("ig") or {}).get("post_count") if last_summary else None,
    )
    _add_top5_slide(prs, "Instagram", ig_posts, "ig")

    _add_plans_slide(prs, slides_data or {})
    _add_closing_slide(prs)

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()
